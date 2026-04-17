from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


OUT_PATH = Path("docs/slides/dataforge_paper_roadmap.pptx")

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(246, 244, 238)
TEXT = RGBColor(24, 30, 36)
MUTED = RGBColor(94, 103, 111)
ACCENT = RGBColor(187, 72, 40)
ACCENT_2 = RGBColor(30, 87, 153)
ACCENT_3 = RGBColor(48, 128, 104)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(219, 214, 206)
WARN = RGBColor(182, 105, 0)

FONT_CN = "Arial Unicode MS"
FONT_EN = "Avenir Next"


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _set_run_font(run, size: int, color: RGBColor, bold: bool = False,
                  font_name: str = FONT_CN, lang: str = "zh-CN") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name

    r_pr = run._r.get_or_add_rPr()
    r_pr.set(qn("a:lang"), lang)

    for child in list(r_pr):
        if child.tag in {qn("a:latin"), qn("a:ea"), qn("a:cs")}:
            r_pr.remove(child)

    r_pr.append(parse_xml(f'<a:latin xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font_name}"/>'))
    r_pr.append(parse_xml(f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font_name}"/>'))
    r_pr.append(parse_xml(f'<a:cs xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font_name}"/>'))


def add_textbox(slide, x, y, w, h, text="", size=20, color=TEXT, bold=False,
                font_name=FONT_CN, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run_font(run, size=size, color=color, bold=bold, font_name=font_name)
    p.alignment = align
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.72, 0.42, 7.8, 0.55, title, size=28, bold=True)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.72), Inches(1.02), Inches(1.2), Inches(0.06)
    ).fill.solid()
    bar = slide.shapes[-1]
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.72, 1.12, 10.8, 0.4, subtitle, size=11, color=MUTED)


def add_bullets(slide, x, y, w, h, items, size=18, color=TEXT, level_indent=0.26):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            _set_run_font(run, size=size, color=color, bold=False, font_name=FONT_CN)
        if level > 0:
            p.left_margin = Inches(level * level_indent)
    return box


def add_card(slide, x, y, w, h, title, lines, accent=ACCENT, title_size=18, body_size=11):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.0)

    strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(0.14), Inches(h)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()

    add_textbox(slide, x + 0.28, y + 0.18, w - 0.42, 0.28, title, size=title_size, bold=True)
    add_bullets(slide, x + 0.28, y + 0.58, w - 0.4, h - 0.7, lines, size=body_size, color=MUTED)


def add_footer(slide, text: str) -> None:
    add_textbox(slide, 0.72, 7.05, 12.0, 0.18, text, size=9, color=MUTED)


def add_table_slide(slide, x, y, w, h, cols, rows):
    table = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for i, col in enumerate(cols):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_2
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                _set_run_font(run, size=12, color=RGBColor(255, 255, 255), bold=True, font_name=FONT_CN)
            p.alignment = PP_ALIGN.CENTER
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if r_idx % 2 == 1 else RGBColor(251, 249, 245)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_run_font(run, size=11, color=TEXT, bold=False, font_name=FONT_CN)
                p.alignment = PP_ALIGN.CENTER if c_idx >= 2 else PP_ALIGN.LEFT
    return table


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    ).fill.solid()
    bg_shape = slide.shapes[-1]
    bg_shape.fill.fore_color.rgb = RGBColor(34, 41, 46)
    bg_shape.line.fill.background()
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(5.9), Inches(13.333), Inches(1.6)
    ).fill.solid()
    band = slide.shapes[-1]
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    add_textbox(slide, 0.8, 0.9, 10.8, 0.7, "DataForge 候选题论文路线图", size=30,
                color=RGBColor(255, 255, 255), bold=True)
    add_textbox(slide, 0.82, 1.7, 10.6, 0.7,
                "围绕系统底座、质量评估/benchmark 与数据方法，规划可持续发表与商业化路线",
                size=16, color=RGBColor(221, 225, 229))
    add_card(
        slide, 0.82, 3.0, 3.55, 1.65, "母体平台", ["DataForge 系统论文", "高并发、限流、恢复、质量门控"], accent=ACCENT_2
    )
    add_card(
        slide, 4.6, 3.0, 3.55, 1.65, "质量平台", ["Quality Report", "Benchmark、Regression"], accent=ACCENT_3
    )
    add_card(
        slide, 8.38, 3.0, 3.55, 1.65, "方法论文", ["代码闭环、RAG 多跳 QA", "Hard negatives、推理轨迹"], accent=WARN
    )
    add_textbox(slide, 0.82, 6.25, 8.0, 0.3, "版本：2026-03-08", size=10, color=RGBColor(255, 245, 240))

    # Strategy slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "总体策略", "不要把每个 pipeline 都做成孤立论文，而是做一个母体平台 + 一个质量平台 + 两条高价值子线。")
    add_card(
        slide, 0.75, 1.6, 3.8, 2.15, "论文分层",
        ["系统论文：DataForge 作为 LLM 数据合成基础设施",
         "方法论文：某类数据如何自动构造、筛选、验证",
         "闭环论文：这些数据是否提升模型或检索器"], accent=ACCENT
    )
    add_card(
        slide, 4.75, 1.6, 3.8, 2.15, "筛选标准",
        ["有明确痛点，而不是单纯功能堆砌",
         "有外部验证信号，而不只是 LLM 自评",
         "能沉淀成基准、数据集或训练 recipe"], accent=ACCENT_2
    )
    add_card(
        slide, 8.75, 1.6, 3.8, 2.15, "推荐主线",
        ["先发 DataForge 系统论文",
         "尽快补 Quality & Benchmark 平台",
         "再选 2 个最硬方向做方法论文"], accent=ACCENT_3
    )
    add_bullets(slide, 0.85, 4.3, 11.5, 1.9, [
        "最值得优先做的组合：系统论文 + 代码执行反馈闭环，或 系统论文 + 企业文档多跳 QA",
        "如果考虑商业化：系统底座之后，优先补 Quality Report、Benchmark Runner、Regression Compare",
        "如果要走长期路线：把推理、偏好、RAG、结构化都做成 DataForge 的领域插件库",
    ], size=18)
    add_footer(slide, "建议形成“母论文 + 子论文 + 闭环论文”的系列化路线。")

    # Quality platform overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "为什么必须补 Quality & Benchmark", "商业化时，客户买的不是 generator 本身，而是可测量的数据质量、收益和可追责性。")
    add_card(
        slide, 0.75, 1.55, 3.8, 2.25, "没有这条线会发生什么",
        ["pipeline 越做越多，但无法比较谁更好",
         "客户只能看到 demo，无法判断值不值得训练/上线",
         "prompt、模型、阈值升级后，无法判断是否回退"], accent=ACCENT
    )
    add_card(
        slide, 4.75, 1.55, 3.8, 2.25, "这条线解决什么",
        ["把“能生成”升级为“能评估、能对比、能回归、能出报告”",
         "支持质量-成本-吞吐-收益的统一比较",
         "形成商业化可签单的质量证明体系"], accent=ACCENT_2
    )
    add_card(
        slide, 8.75, 1.55, 3.8, 2.25, "最关键的输出物",
        ["数据质量报告",
         "benchmark 对比报告",
         "版本回归变化报告"], accent=ACCENT_3
    )
    add_bullets(slide, 0.85, 4.25, 11.7, 1.6, [
        "研究上，它会把 DataForge 从“系统工具”升级成“数据生产 + 质量评测”平台",
        "产品上，它会让客户看到标准化报告，而不是只看到 jsonl 文件",
        "因此 Quality & Benchmark 的优先级不低于继续新增 3-5 个 pipeline"
    ], size=18)
    add_footer(slide, "结论：这不是附属功能，而是商业化前必须补的第二条主轴。")

    # Quality architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Quality & Benchmark 产品结构", "建议拆成 4 层：单条样本评估、整批数据报告、benchmark runner、回归治理。")
    add_card(slide, 0.75, 1.52, 2.85, 2.4, "1. Record-level", [
        "格式、内容、推理、代码、RAG、安全 probes",
        "输出 RecordEvaluation"
    ], accent=ACCENT)
    add_card(slide, 3.9, 1.52, 2.85, 2.4, "2. Dataset-level", [
        "acceptance、score、重复率、风险、成本",
        "输出 DatasetQualityReport"
    ], accent=ACCENT_2)
    add_card(slide, 7.05, 1.52, 2.85, 2.4, "3. Benchmark", [
        "Reasoning / Code / RAG mini benches",
        "输出 BenchmarkResult"
    ], accent=ACCENT_3)
    add_card(slide, 10.2, 1.52, 2.35, 2.4, "4. Regression", [
        "比较新旧 run",
        "输出 RegressionReport"
    ], accent=WARN)
    add_bullets(slide, 0.85, 4.3, 11.75, 1.55, [
        "短期先做 report + compare；中期再做 mini benchmark；长期再做行业特化质量标准",
        "核心不是某一个 judge，而是统一 schema、统一报告、统一回归口径",
        "这条线既能服务论文，也能直接转成商业产品能力"
    ], size=18)
    add_footer(slide, "推荐最先落地的 3 个模块：DatasetReport、RegressionCompare、BenchmarkRunner。")

    # Quality MVP
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Quality & Benchmark MVP", "8 周内可落地的版本：先让每次 pipeline 运行后都自动生成标准化质量报告。")
    add_card(slide, 0.75, 1.48, 3.8, 2.55, "MVP 模块", [
        "Record Evaluator Adapter",
        "Dataset Report Builder",
        "Regression Compare",
        "Minimal Benchmark Runner"
    ], accent=ACCENT)
    add_card(slide, 4.78, 1.48, 3.8, 2.55, "统一对象", [
        "RecordEvaluation",
        "DatasetQualityReport",
        "BenchmarkResult",
        "RegressionReport"
    ], accent=ACCENT_2)
    add_card(slide, 8.81, 1.48, 3.8, 2.55, "CLI", [
        "dataforge report run_output.jsonl",
        "dataforge compare report_a.json report_b.json",
        "dataforge benchmark run config.yaml --bench code-mini"
    ], accent=ACCENT_3)
    add_bullets(slide, 0.85, 4.35, 11.7, 1.55, [
        "MVP 指标先只做 12 个：acceptance、score 分位数、schema valid、duplicate、PII、cost、throughput 等",
        "MVP benchmark 先只做 3 类：ReasoningMiniBench、CodeMiniBench、RAGMiniBench",
        "MVP 报告同时导出 report.json、report.md、record_evaluations.jsonl"
    ], size=18)
    add_footer(slide, "优先级：和新增方法论文同级，甚至更高。")

    # Candidate matrix
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "六个候选题总览", "在系统主线之外，再加一条 Quality & Benchmark 主线，形成研究与商业的双轮驱动。")
    rows = [
        ["DataForge 系统论文", "MLSys / ACL Industry / Demo", "S", "中", "最快可起稿"],
        ["代码执行反馈闭环", "ACL / EMNLP / ICSE", "S", "中高", "最硬、最可验证"],
        ["企业文档多跳 QA", "ACL / EMNLP / SIGIR", "S", "中高", "产业壁垒强"],
        ["Hard Negatives", "SIGIR / ACL", "A", "中", "低算力、见效快"],
        ["推理轨迹质量控制", "ACL / EMNLP", "A", "高", "竞争激烈"],
        ["企业偏好与宪法对齐", "ACL Industry / EMNLP", "A-", "中", "企业价值高"],
    ]
    add_table_slide(slide, 0.72, 1.45, 11.9, 4.6,
                    ["候选题", "推荐 venue", "优先级", "研究风险", "一句话判断"], rows)
    add_bullets(slide, 0.82, 6.25, 12.0, 0.55, [
        "补充判断：所有候选题都应优先接入统一质量报告和 benchmark，否则很难形成可比较的产品线"
    ], size=14, color=WARN)

    # Systems paper
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "候选题 1：DataForge 系统论文",
              "定位：面向 LLM 数据合成的高并发、可恢复、质量可控基础设施。")
    add_card(slide, 0.75, 1.5, 3.8, 2.3, "论文题目草案", [
        "DataForge: A Fault-Tolerant and Rate-Aware Infrastructure for LLM Data Synthesis",
        "关键词：dual-bucket、WAL、async pipeline、quality gates"
    ], accent=ACCENT)
    add_card(slide, 4.75, 1.5, 3.8, 2.3, "核心创新点", [
        "把并发执行、双维限流、断点恢复、评估链路统一到一个数据生产系统",
        "不是单次 prompt，而是长期运行的训练数据生产工作流"
    ], accent=ACCENT_2)
    add_card(slide, 8.75, 1.5, 3.8, 2.3, "最小实验包", [
        "吞吐：sequential / naive async / threaded / DataForge",
        "限流：dual bucket / rpm-only / no limiter",
        "恢复：no checkpoint / JSONL WAL / SQLite WAL"
    ], accent=ACCENT_3)
    add_bullets(slide, 0.85, 4.15, 11.8, 1.7, [
        "你现在仓库里已经有大部分实验基础，可先作为整个研究计划的第一篇",
        "风险不在实现，而在 novelty 表述：不能只说“做了个框架”，必须强调 LLM 数据生产的系统问题",
        "推荐 venue：MLSys；若补企业实践和产品化，也可考虑 ACL/NAACL/EMNLP Industry 或 Demo"
    ], size=18)
    add_footer(slide, "优先级：S。最适合作为后续所有子论文的母体平台。")

    # Code loop
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "候选题 2：代码执行反馈闭环数据",
              "定位：从“生成代码”升级到“失败-诊断-修复”的可验证训练数据。")
    add_card(slide, 0.75, 1.5, 4.0, 2.3, "为什么最值得做", [
        "外部执行器提供硬反馈，不容易被 reviewer 认为只是 prompt engineering",
        "天然适合和 DataForge 的并发、重试、沙盒、失败恢复绑定"
    ], accent=ACCENT)
    add_card(slide, 4.95, 1.5, 4.0, 2.3, "最小实验包", [
        "question -> code -> execution feedback -> repair 数据链路",
        "compile rate / unit test pass rate / pass@1 / 修复轮数 / token 成本"
    ], accent=ACCENT_2)
    add_card(slide, 9.15, 1.5, 3.45, 2.3, "推荐 venue", [
        "ACL / EMNLP：更偏数据方法和模型提升",
        "ICSE：更偏软件工程工作流与代码修复"
    ], accent=ACCENT_3)
    add_bullets(slide, 0.85, 4.15, 11.9, 1.7, [
        "关键不是“模型会不会写代码”，而是“反馈闭环数据是否显著提升修复与调试能力”",
        "建议先做小规模 benchmark 和稳定沙盒，再扩大数据量，避免大量失败样本浪费算力",
        "这是最可能做成强 paper 的子方向之一"
    ], size=18)
    add_footer(slide, "优先级：S。推荐作为系统论文后的第一条方法线。")

    # RAG + hard negatives
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "候选题 3-4：企业文档多跳 QA 与 Hard Negatives",
              "定位：围绕企业知识库，把‘死文档’转成高价值检索与 RAG 训练数据。")
    add_card(slide, 0.75, 1.52, 5.8, 2.5, "Doc2QA / Multi-hop QA", [
        "从 PDF、研报、制度文档中自动构造问题、答案、证据链",
        "重点不在问答本身，而在跨 chunk、跨页面、跨文档的证据拼接",
        "指标：answer accuracy、evidence recall、citation faithfulness"
    ], accent=ACCENT_2)
    add_card(slide, 6.78, 1.52, 5.8, 2.5, "Hard Negatives", [
        "生成“词面相关但证据错误”的高难负样本，服务 embedding / reranker",
        "指标：MRR、NDCG、Recall@k、域外迁移",
        "优点：低算力、见效快、最适合先打检索实验"
    ], accent=ACCENT_3)
    add_bullets(slide, 0.85, 4.4, 11.8, 1.6, [
        "如果你更看重行业壁垒，优先做多跳 QA；如果你更看重低成本快速发文，优先做 Hard Negatives",
        "这条线特别适合金融、医疗、法务、制造等私有文档场景",
        "推荐 venue：ACL / EMNLP / SIGIR；如果案例够真实，industry track 也很合适"
    ], size=18)
    add_footer(slide, "优先级：Multi-hop QA = S，Hard Negatives = A。")

    # Reasoning traces
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "候选题 5：推理轨迹质量控制",
              "定位：不是只生成 CoT，而是研究哪些 trace 真能提升小模型。")
    add_card(slide, 0.75, 1.5, 3.9, 2.2, "研究问题", [
        "哪些 synthetic reasoning traces 有效，哪些只是冗长噪声",
        "是否需要 self-correction，以及插入在哪一步最有效"
    ], accent=ACCENT)
    add_card(slide, 4.85, 1.5, 3.9, 2.2, "实验设计", [
        "direct answer / naive CoT / self-correction CoT / filtered CoT",
        "评测 GSM8K、MATH 子集或垂域推理 benchmark"
    ], accent=ACCENT_2)
    add_card(slide, 8.95, 1.5, 3.6, 2.2, "风险", [
        "社区竞争最激烈",
        "必须把“质量控制”讲清楚，否则容易沦为重复工作"
    ], accent=WARN)
    add_bullets(slide, 0.85, 4.15, 11.8, 1.8, [
        "推荐做法：把多维评估、trace 过滤、长度控制和蒸馏收益绑定起来，而不是只展示更长的 <think> 输出",
        "如果没有下游小模型收益，这个方向的说服力会明显下降",
        "适合在系统论文之后，作为更偏 NLP 的方法线推进"
    ], size=18)
    add_footer(slide, "优先级：A。值得做，但不建议作为第一条子线。")

    # Preference and structured
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "候选题 6：企业偏好 / 宪法对齐；结构化与脱敏",
              "定位：偏行业价值强，适合 industry track 或垂域论文。")
    add_card(slide, 0.75, 1.48, 5.8, 2.55, "企业偏好 / Constitutional Alignment", [
        "把企业规则转成 chosen/rejected 或 rewrite 轨迹，而不是依赖全人工标注",
        "适合风格、合规、安全、简洁、客服 SOP 等场景",
        "关键是处理 judge 偏差和 helpfulness 损失"
    ], accent=ACCENT)
    add_card(slide, 6.78, 1.48, 5.8, 2.55, "结构化 / 脱敏数据", [
        "重点不只是 schema 合法，而是跨字段逻辑、业务规则、repair loop 和 privacy-utility 平衡",
        "更像行业论文或系统能力扩展，不建议单独从‘格式校验’立题"
    ], accent=ACCENT_3)
    add_bullets(slide, 0.85, 4.38, 11.8, 1.55, [
        "这两条都很适合做成 DataForge 的企业插件线，但学术 novelty 要靠真实规则、真实基准或下游收益撑起来",
        "如果数据不能公开，industry track 往往比常规 research track 更合适",
        "建议放在系统论文和核心方法论文之后推进"
    ], size=18)
    add_footer(slide, "优先级：A- 到 B+。更偏产业落地与后续扩展。")

    # Compute budget
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "算力预算估计", "按 DataForge 当前 vLLM 实验吞吐和 LoRA 小模型训练口径做的保守估算。")
    rows = [
        ["系统论文", "10-26 GPUh", "0", "2-6 GPUh", "10-26 GPUh"],
        ["推理轨迹", "10-30 GPUh", "12-30 GPUh / 50-120 GPUh", "4-10 GPUh", "26-70 / 64-160 GPUh"],
        ["代码闭环", "15-40 GPUh", "10-25 GPUh / 40-100 GPUh", "5-15 GPUh", "30-80 / 60-155 GPUh"],
        ["偏好对齐", "8-20 GPUh", "8-20 GPUh / 30-80 GPUh", "4-8 GPUh", "20-48 / 42-108 GPUh"],
        ["多跳 QA", "10-25 GPUh", "10-25 GPUh / 40-100 GPUh", "4-10 GPUh", "24-60 / 54-135 GPUh"],
        ["Hard Negatives", "6-15 GPUh", "6-16 GPUh / 12-24 GPUh", "2-6 GPUh", "14-37 GPUh"],
    ]
    add_table_slide(slide, 0.72, 1.4, 11.9, 4.6,
                    ["方向", "合成", "训练（小模型 / 7B）", "评测", "总计"], rows)
    add_bullets(slide, 0.82, 6.22, 12.0, 0.55, [
        "低预算路线：系统论文 + Hard Negatives；平衡路线：系统论文 + 多跳 QA；高价值路线：系统论文 + 代码闭环"
    ], size=13, color=WARN)

    # Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "18 个月路线图", "以系统底座为中心，先补质量平台，再扩展强方法线和收益闭环。")
    add_card(slide, 0.75, 1.55, 2.8, 2.2, "阶段 1", [
        "3-4 个月",
        "完成 DataForge 系统论文",
        "补齐 throughput / rate-limit / recovery 实验"
    ], accent=ACCENT)
    add_card(slide, 3.75, 1.55, 2.8, 2.2, "阶段 2", [
        "2-3 个月",
        "补 Quality & Benchmark MVP",
        "先上线 report / compare / mini bench"
    ], accent=ACCENT_2)
    add_card(slide, 6.75, 1.55, 2.8, 2.2, "阶段 3", [
        "4-6 个月",
        "选一条硬方法线",
        "优先：代码闭环 或 多跳 QA"
    ], accent=ACCENT_3)
    add_card(slide, 9.75, 1.55, 2.8, 2.2, "阶段 4", [
        "4-6 个月",
        "第二条子线或闭环论文",
        "回答：这些数据到底提升了什么"
    ], accent=WARN)
    add_bullets(slide, 0.85, 4.3, 11.8, 1.7, [
        "最稳路线：DataForge 系统论文 -> Quality & Benchmark MVP -> 代码执行反馈闭环 / 企业文档多跳 QA",
        "最省算力路线：DataForge 系统论文 -> Quality MVP -> Hard Negatives -> 企业偏好 / 结构化行业论文",
        "核心原则：每条新 pipeline 都优先复用母体平台与统一评测体系，不要散成独立脚本"
    ], size=18)
    add_footer(slide, "推荐以 2 条主方法线为核心，不建议同时铺开 5-6 条。")

    # Final recommendation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "最终建议", "先收敛，再扩展；先补质量平台，再铺开更多 pipeline。")
    add_card(slide, 0.75, 1.48, 3.8, 2.45, "第一优先级", [
        "DataForge 系统论文",
        "最快落地，且是后续所有方法线的母体"
    ], accent=ACCENT)
    add_card(slide, 4.78, 1.48, 3.8, 2.45, "第二优先级", [
        "Quality & Benchmark MVP",
        "决定商业化说服力，也决定各 pipeline 是否可比较"
    ], accent=ACCENT_2)
    add_card(slide, 8.81, 1.48, 3.8, 2.45, "第三优先级", [
        "代码闭环 / 多跳 QA / Hard Negatives",
        "按算力与目标选择其中 1-2 条作为主方法线"
    ], accent=ACCENT_3)
    add_bullets(slide, 0.85, 4.38, 11.7, 1.45, [
        "如果近期资源有限，先做“系统论文 + Quality MVP + Hard Negatives”",
        "如果目标是冲更强论文，做“系统论文 + Quality MVP + 代码闭环”",
        "如果目标是产业价值和垂域壁垒，做“系统论文 + Quality MVP + 企业文档多跳 QA”"
    ], size=19)
    add_textbox(slide, 0.85, 6.15, 10.8, 0.42,
                "一句话：先把 DataForge 和 Quality 平台立起来，再让最硬的两个数据方向替平台证明价值。",
                size=19, color=ACCENT, bold=True)

    return prs


def main() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT_PATH)
    print(f"Generated: {OUT_PATH}")


if __name__ == "__main__":
    main()
